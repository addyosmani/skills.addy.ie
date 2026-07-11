#!/usr/bin/env python3
"""Build 101/201/301 PPTX decks from scripts/deck-data.json.

Dark, on-brand slides that mirror the HTML decks. Diagram slides embed the
pre-rendered PNGs from public/teach/diagrams/. No em dashes.

    python3 scripts/build_pptx.py
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(ROOT, "scripts", "deck-data.json")
DIAGRAMS = os.path.join(ROOT, "public", "teach", "diagrams")
OUT = os.path.join(ROOT, "public", "teach", "decks")

# Palette
BG = RGBColor(0x0A, 0x0A, 0x0A)
CARD = RGBColor(0x12, 0x12, 0x12)
BORDER = RGBColor(0x2E, 0x2E, 0x2E)
FG = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0xED, 0xED, 0xED)
MUTED = RGBColor(0xA1, 0xA1, 0xA1)
DIM = RGBColor(0x7D, 0x7D, 0x7D)
FAINT = RGBColor(0x56, 0x56, 0x56)

SANS = "Geist"
MONO = "Geist Mono"

EMU_IN = 914400
W = Inches(13.333)
H = Inches(7.5)


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, font=SANS, first=False,
         align=PP_ALIGN.LEFT, space_after=0, space_before=0, spacing=None,
         tracking=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    if space_before:
        p.space_before = Pt(space_before)
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font
    f.color.rgb = color
    return p


def kicker(slide, text, accent, x=Inches(0.9), y=Inches(0.85)):
    tf = textbox(slide, x, y, Inches(10), Inches(0.4))
    p = para(tf, text.upper(), 13, accent, bold=True, font=MONO, first=True)


def footer(slide, deck, idx, total):
    # left brand
    tf = textbox(slide, Inches(0.9), Inches(6.95), Inches(4), Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "agent"; r.font.name = MONO; r.font.size = Pt(11); r.font.color.rgb = MUTED
    r2 = p.add_run(); r2.text = "-skills"; r2.font.name = MONO; r2.font.size = Pt(11); r2.font.color.rgb = FAINT
    # middle
    tf2 = textbox(slide, Inches(4.66), Inches(6.95), Inches(4), Inches(0.4))
    para(tf2, deck["title"], 11, FAINT, font=MONO, first=True, align=PP_ALIGN.CENTER)
    # page
    tf3 = textbox(slide, Inches(9.43), Inches(6.95), Inches(3), Inches(0.4))
    para(tf3, f"{idx} / {total}", 11, FAINT, font=MONO, first=True, align=PP_ALIGN.RIGHT)


def cmd_box(slide, code, accent, y):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(8.4), Inches(0.95))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0E, 0x0E, 0x0E)
    box.line.color.rgb = BORDER; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.28); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "$ "; r.font.name = MONO; r.font.size = Pt(20); r.font.color.rgb = accent
    r2 = p.add_run(); r2.text = code; r2.font.name = MONO; r2.font.size = Pt(20); r2.font.color.rgb = BODY


def render(deck):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    accent = hex_to_rgb(deck["accent"])
    total = len(deck["slides"])

    for i, s in enumerate(deck["slides"], start=1):
        slide = add_slide(prs)
        t = s["type"]

        if t == "title":
            kicker(slide, s["kicker"], accent)
            tf = textbox(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2))
            para(tf, s["title"], 66, FG, bold=True, first=True, spacing=1.0)
            tf2 = textbox(slide, Inches(0.9), Inches(4.2), Inches(9.5), Inches(1.3))
            para(tf2, s["subtitle"], 24, MUTED, first=True, spacing=1.1)
            tf3 = textbox(slide, Inches(0.9), Inches(5.7), Inches(9), Inches(0.6))
            para(tf3, s["foot"], 16, DIM, first=True)

        elif t == "statement":
            kicker(slide, s["kicker"], accent)
            tf = textbox(slide, Inches(0.9), Inches(2.0), Inches(11), Inches(2.4))
            para(tf, s["text"], 52, FG, bold=True, first=True, spacing=1.02)
            tf2 = textbox(slide, Inches(0.9), Inches(4.7), Inches(10), Inches(1.6))
            para(tf2, s["sub"], 22, MUTED, first=True, spacing=1.2)

        elif t == "bullets":
            kicker(slide, s["kicker"], accent)
            tf = textbox(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.2))
            para(tf, s["heading"], 40, FG, bold=True, first=True, spacing=1.05)
            tfb = textbox(slide, Inches(0.9), Inches(3.0), Inches(11.4), Inches(3.6))
            for j, b in enumerate(s["bullets"]):
                p = tfb.paragraphs[0] if j == 0 else tfb.add_paragraph()
                p.space_after = Pt(14)
                p.line_spacing = 1.15
                marker = p.add_run()
                marker.text = "-  "
                marker.font.name = MONO
                marker.font.size = Pt(22)
                marker.font.color.rgb = accent
                r = p.add_run()
                r.text = b
                r.font.name = SANS
                r.font.size = Pt(22)
                r.font.color.rgb = BODY

        elif t == "diagram":
            kicker(slide, s["kicker"], accent)
            tf = textbox(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.9))
            para(tf, s["heading"], 34, FG, bold=True, first=True)
            name = s["image"].split("/")[-1].replace(".svg", ".png")
            png = os.path.join(DIAGRAMS, name)
            if os.path.exists(png):
                # 16:9 image, fit width ~ 9.7in centered-left
                pic_w = Inches(9.9)
                pic_h = Emu(int(pic_w * 9 / 16))
                left = Inches(0.9)
                top = Inches(2.35)
                slide.shapes.add_picture(png, left, top, width=pic_w, height=pic_h)
            tfc = textbox(slide, Inches(0.9), Inches(6.35), Inches(11.4), Inches(0.6))
            para(tfc, s["caption"], 15, DIM, first=True)

        elif t == "code":
            kicker(slide, s["kicker"], accent)
            tf = textbox(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.2))
            para(tf, s["heading"], 40, FG, bold=True, first=True, spacing=1.05)
            cmd_box(slide, s["code"], accent, Inches(3.2))
            tf2 = textbox(slide, Inches(0.9), Inches(4.5), Inches(10.5), Inches(1.8))
            para(tf2, s["sub"], 21, MUTED, first=True, spacing=1.2)

        elif t == "commands":
            kicker(slide, s["kicker"], accent)
            tf = textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.9))
            para(tf, s["heading"], 34, FG, bold=True, first=True)
            rows = s["rows"]
            table = slide.shapes.add_table(len(rows), 3, Inches(0.9), Inches(2.35),
                                           Inches(11.5), Inches(0.5 * len(rows))).table
            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(5.0)
            table.columns[2].width = Inches(3.5)
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    cell = table.cell(ri, ci)
                    cell.fill.solid(); cell.fill.fore_color.rgb = BG
                    cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
                    cell.margin_left = Pt(6)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = cell.text_frame.paragraphs[0]
                    r = p.add_run(); r.text = val
                    r.font.size = Pt(16)
                    if ci == 0:
                        r.font.name = MONO; r.font.color.rgb = accent
                    elif ci == 1:
                        r.font.name = SANS; r.font.color.rgb = BODY
                    else:
                        r.font.name = SANS; r.font.color.rgb = FAINT; p.alignment = PP_ALIGN.RIGHT

        elif t == "cta":
            tf = textbox(slide, Inches(1.2), Inches(2.0), Inches(10.9), Inches(2), anchor=MSO_ANCHOR.MIDDLE)
            para(tf, s["title"], 44, FG, bold=True, first=True, align=PP_ALIGN.CENTER, spacing=1.05)
            cmd_box(slide, s["code"], accent, Inches(4.2))
            tf2 = textbox(slide, Inches(1.2), Inches(5.5), Inches(10.9), Inches(0.6))
            para(tf2, "   ".join(s["links"]), 16, DIM, first=True, font=MONO, align=PP_ALIGN.CENTER)

        footer(slide, deck, i, total)

    out = os.path.join(OUT, deck["slug"] + ".pptx")
    prs.save(out)
    print("wrote", os.path.relpath(out, ROOT), f"({total} slides)")


def main():
    with open(DATA) as f:
        decks = json.load(f)
    os.makedirs(OUT, exist_ok=True)
    for deck in decks:
        render(deck)


if __name__ == "__main__":
    main()
